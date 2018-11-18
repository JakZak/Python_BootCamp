import csv

from pptx import Presentation

prs = Presentation()

slide_layout = prs.slide_layouts[1]

slide = prs.slides.add_slide(slide_layout)

shapes = slide.shapes

titla_shape = shapes.title

body_shape = shapes.placeholders[1]

tittle_shape = "Jakiś tekst"

tf = body_shape.text_frame
tf.text = "Zawartość text frame"

with open ()
p = tf.add_paragraph()
p.text = "Kobiety"
p.level = 1

p = tf.add_paragraph()
p.text = "Przeżyło: 70%"
p.level = 2

prs.save("raport_titanic.pptx")